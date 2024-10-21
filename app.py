import streamlit as st
from dotenv import load_dotenv
import os
import google.generativeai as genai
from youtube_transcript_api import YouTubeTranscriptApi
from PyPDF2 import PdfReader
from gtts import gTTS
import tempfile
import re
# Load environment variables
load_dotenv()

# Configure Google Gemini API
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Function to generate a summary using Google Gemini Pro
def generate_summary(text, prompt):
    model = genai.GenerativeModel("gemini-pro")
    response = model.generate_content(prompt + text)
    return response.text

# Function to extract YouTube video transcript
def extract_transcript_details(youtube_video_url):
    try:
        video_id = youtube_video_url.split("=")[1]
        transcript_text = YouTubeTranscriptApi.get_transcript(video_id)
        transcript = " ".join([i["text"] for i in transcript_text])
        return transcript
    except Exception as e:
        raise e

# Function to convert text to speech using Google Text-to-Speech (gTTS)
def text_to_speech(text):
    clean_text = re.sub(r"[^\w\s,.!?']", "", text)  # Clean the text from special characters
    tts = gTTS(text=clean_text, lang='en')  # Convert text to speech
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')  # Temporary audio file
    tts.save(temp_file.name)  # Save audio
    return temp_file.name  # Return file path


# Function to generate the summarization prompt based on the selected style
def get_summarization_prompt(style):
    if style == "Bullet Points":
        return "Summarize the video in clear bullet points (max 250 words): "
    elif style == "Detailed Paragraph":
        return "Provide a detailed paragraph summarizing the video (max 250 words): "
    elif style == "Short Summary":
        return "Summarize the video briefly (max 100 words): "
    elif style == "Key Highlights":
        return "Highlight key points and action items from the video (max 150 words): "

# Function to extract the YouTube video ID from various URL formats
def extract_video_id(youtube_video_url):
    pattern = r"(?:v=|\/)([0-9A-Za-z_-]{11}).*"
    match = re.search(pattern, youtube_video_url)
    if match:
        return match.group(1)
    else:
        raise ValueError("Invalid YouTube URL format")
    
# Set page configuration
st.set_page_config(page_title="AI App Suite", layout="wide")

# Sidebar navigation with all pages
st.sidebar.image("./images/Youtube_Logo.png", width=150)
st.sidebar.title("AI App Suite")

# Add all pages to the navigation
# nav_option = st.sidebar.radio("Navigation", ["About Us", "Summarize YouTube Videos", "PDF Summarizer", "Chat to PDF", "PPT Maker"])
nav_option = st.sidebar.radio(
    "Navigation", 
    [
        "📄 About Us", 
        "🎥 Summarize YouTube Videos", 
        "📑 PDF Summarizer", 
        "📜 Chat to PDF", 
        "📊 PPT Maker"
    ]
)


st.markdown("""<style>
    .step-container {
        background-color: #f9f9f9;
        padding: 20px;
        border-radius: 10px;
        border: 2px solid #e6e6e6;
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);
        max-width: 800px;
        margin: 20px auto;
        text-align: left;
        font-family: 'Arial', sans-serif;
    }
    
    .step-container h3 {
        color: #3a7ca5;
        font-size: 24px;
    }

    .step-container p {
        color: #555555;
        font-size: 18px;
        line-height: 1.6;
    }

    .step-container strong {
        color: #2a5478;
    }</style>""", unsafe_allow_html=True)



# Summarize YouTube Videos Page
if nav_option == "🎥 Summarize YouTube Videos":
    st.title("YouTube Video Summary Tool")
    youtube_link = st.text_input("Enter YouTube Video Link:", placeholder="Paste the link here")

    if youtube_link:
        video_id = youtube_link.split("=")[1]
        st.image(f"http://img.youtube.com/vi/{video_id}/0.jpg", use_column_width=True)

    # Summarization style options
    summary_style = st.selectbox(
        "Select Summarization Style", 
        ["Bullet Points", "Detailed Paragraph", "Short Summary", "Key Highlights"]
    )

    if st.button("Get Detailed Notes"):
        try:
            transcript_text = extract_transcript_details(youtube_link)
            if transcript_text:
                summary_prompt = get_summarization_prompt(summary_style)
                summary = generate_summary(transcript_text, summary_prompt)
                st.markdown("## Detailed Notes:")
                st.write(summary)
                st.session_state.summary = summary
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")
    if 'summary' in st.session_state:
        st.markdown("### Summary Text:")
        st.write(st.session_state.summary)

        if st.button("Listen to Summary"):
            summary_text = st.session_state.summary
            audio_file_path = text_to_speech(summary_text)
            st.audio(audio_file_path, format='audio/mp3')
    
    st.title("How to Summarize YouTube Videos?")
    

# Applying the CSS styles to the content
    st.markdown("""<div class="step-container">
        <h3>Follow These 3 Simple Steps to Summarize a YouTube Video:</h3>
        <p><strong>Step 1:</strong> Get the YouTube video link<br>
        Copy and paste the YouTube video link into NoteGPT.</p>
        <p><strong>Step 2:</strong> Generate a Summary<br>
        Click the "Generate Summary" button, and NoteGPT will fetch the transcript and summarize the YouTube video.</p>
        <p><strong>Step 3:</strong> Read the AI summary<br>
        Read the concise summary and save valuable time.</p>
    </div>""", unsafe_allow_html=True)


# PDF Summarizer Page
elif nav_option == "📑 PDF Summarizer":
    st.title("AI PDF Summarizer")
    st.markdown("""
    AI PDF Summarizer can summarize long PDF documents in seconds. 
    It can convert PDFs to text, allowing you to ask your PDF.
    The best PDF Summary tool helps you save time and learning faster and better.
    """)
    # File uploader for PDF files
    uploaded_file = st.file_uploader("Upload PDF", type="pdf")

    if uploaded_file:
        reader = PdfReader(uploaded_file)
        text = ""

        # Extract text from each page
        for page in reader.pages:
            text += page.extract_text()

        if text:
            st.text_area("Extracted Text", value=text, height=300)
            
            # Button to summarize the extracted text
            if st.button("Summarize PDF"):
                prompt = """You are an AI text summarizer. Please summarize the following text 
                            into important points within 250 words. Here is the text: """
                with st.spinner("Generating summary..."):
                    summary = generate_summary(text, prompt)
                st.success("Summary Generated!")
                st.markdown("## Summary:")
                st.write(summary)
        else:
            st.warning("No text found in the PDF file.")

    st.title("Spending too much time on lengthy PDF files?")
    st.markdown("""
    Say goodbye to time-consuming PDF summaries with NoteGPT's PDF Summary tool.
    AI PDF Summarizer is free online tool saves you time and enhances your learning experience.
    The PDF Summarizer can convert PDFs to text page by page to and summarize Large PDFs into concise summaries and PDF to mind map with just one click.
    """)

    st.title("How to Use AI PDF Summarizer")
    st.markdown("""<div class="step-container">
        <h3>How to Use AI PDF Summarizer in 3 Simple Steps:</h3>
        <p><strong>Step 1:</strong> Upload Your PDF<br>
        Drag and drop your PDF file into the AI PDF Summarizer, or enter the PDF URL. You can also upload directly from Google Drive (coming soon).</p>
        <p><strong>Step 2:</strong> Select Your Desired Feature<br>
        Choose from a range of AI-powered features such as summarizing the PDF document, extracting PDF to text, asking your PDF or chatting with AI, translating PDF content, generating mind maps from the PDF, or using the AI PDF Reader.</p>
        <p><strong>Step 3:</strong> Save and Share<br>
        Save the summarized content, extracted text, or generated mind maps. You can also share your results with friends or colleagues for collaborative learning and discussion.</p>
        </div>""", unsafe_allow_html=True)


# Chat to PDF Page
elif nav_option == "📜 Chat to PDF":
    st.title("Chat with your any PDF")
    st.markdown("""
        Chat with any PDF using NoteGPT's ChatPDF app, along with its alternative solutions.
        You can ask your PDF documents questions, seek answers, summaries, and more.
        all free and accessible.
    """)
    chat_input = st.text_area("Enter your text to convert to PDF")

    if st.button("Generate PDF"):
        from fpdf import FPDF
        
        if chat_input:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, chat_input)
            pdf.output("chat_to_pdf.pdf")
            
            with open("chat_to_pdf.pdf", "rb") as file:
                st.download_button(label="Download PDF", data=file, file_name="chat_to_pdf.pdf", mime="application/pdf")
        else:
            st.warning("Please enter text to convert.")
    
    st.title("Can I chat with PDFs using AI?")
    st.markdown("""
        ChatPDF allowing users to chat with PDF documents using AI.
        This tool provides AI-powered assistance, enabling users to ask with PDF for questions, extract insights, and analyze PDF content effortlessly.
        ChatPDF is free and accessible to all users, offering versatile features suitable for students, researchers, and professionals. It provides a convenient and efficient way to interact with PDFs.
    """)

    st.title("How to Use ChatPDF?")
    st.markdown("""<div class="step-container">
        <h3>Ready to Get Started? Unlock Interactive Possibilities with ChatPDF!</h3>
        <p><strong>Step 1:</strong> Upload Your PDF<br>
        Begin by uploading your PDF document to ChatPDF. Simply click on the upload button and select the file you want to chat with.</p>
        <p><strong>Step 2:</strong> Ask with Your PDF<br>
        Once your PDF is uploaded, start asking questions or engaging in conversation with ChatPDF. You can inquire about specific content, seek summaries, or extract insights from the document.</p>
        <p><strong>Step 3:</strong> Explore Insights and Summaries<br>
        ChatPDF will provide you with AI-powered insights, summaries, and answers based on your queries. Explore the extracted information and make the most out of your PDF interaction experience.</p>
    </div>""", unsafe_allow_html=True)


# PPT Maker Page
elif nav_option == "📊 PPT Maker":
    st.title("PPT Maker")

    # Option to either input slide content manually or upload a file
    slide_option = st.radio("Choose how to add content to your slides:", ("Input Manually", "Upload Text File (.txt)"))

    # If user chooses to input slide content manually
    if slide_option == "Input Manually":
        st.subheader("Enter Slide Content:")
        
        # Get title and content for each slide
        slide_title = st.text_input("Enter Slide Title:")
        slide_content = st.text_area("Enter Slide Content:", height=200)

        # Button to generate the PPT
        if st.button("Generate PPT"):
            from pptx import Presentation
            from pptx.util import Inches

            # Create a presentation object
            prs = Presentation()

            # Add a slide with title and content layout
            slide_layout = prs.slide_layouts[1]  # 1 is for a title and content slide
            slide = prs.slides.add_slide(slide_layout)

            # Set slide title and content
            slide.shapes.title.text = slide_title
            slide.shapes.placeholders[1].text = slide_content

            # Save the presentation
            ppt_file = "generated_ppt.pptx"
            prs.save(ppt_file)

            # Provide download button
            with open(ppt_file, "rb") as file:
                st.download_button(
                    label="Download PPT",
                    data=file,
                    file_name="generated_ppt.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
    
    # If user chooses to upload a text file for content
    elif slide_option == "Upload Text File (.txt)":
        uploaded_file = st.file_uploader("Upload a Text File", type="txt")

        if uploaded_file:
            content = uploaded_file.read().decode("utf-8")
            st.text_area("Uploaded Content", content, height=300)

            if st.button("Generate PPT from File"):
                from pptx import Presentation
                from pptx.util import Inches

                # Split content by line for each slide
                lines = content.split("\n")
                
                # Create a presentation object
                prs = Presentation()

                for i in range(0, len(lines), 2):  # Assuming alternate lines for title and content
                    slide_layout = prs.slide_layouts[1]  # 1 is for title and content slide
                    slide = prs.slides.add_slide(slide_layout)

                    # Set slide title and content
                    slide_title = lines[i].strip() if i < len(lines) else "No Title"
                    slide_content = lines[i + 1].strip() if i + 1 < len(lines) else "No Content"
                    slide.shapes.title.text = slide_title
                    slide.shapes.placeholders[1].text = slide_content

                # Save the presentation
                ppt_file = "generated_ppt.pptx"
                prs.save(ppt_file)

                # Provide download button
                with open(ppt_file, "rb") as file:
                    st.download_button(
                        label="Download PPT",
                        data=file,
                        file_name="generated_ppt.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

    else:
        st.warning("Please choose how you'd like to add content to the slides.")


# About Us Page
elif nav_option == "📄 About Us":
    st.title("Welcome to AI App Suite!")
    st.markdown("""
    ### Explore our AI-driven solutions:
    
    - **🎥 Summarize YouTube Videos**: Automatically generate concise summaries from YouTube videos.
    - **📑 PDF Summarizer**: Instantly convert and summarize PDF documents.
    - **📜 Chat to PDF**: Enter text, and instantly convert it into a shareable PDF.
    - **📊 PPT Maker**: Quickly generate professional PowerPoint presentations from your notes.
    
    Leverage the power of AI to save time and simplify your work. With an easy-to-use interface and cutting-edge tools, AI App Suite is your one-stop shop for productivity.
    """)


# Footer
st.sidebar.markdown("---")
st.sidebar.markdown("Copyright © 2024; Designed ❤️ by Code Rangers")
